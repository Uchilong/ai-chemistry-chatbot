import os
from typing import Optional, Dict
import google.generativeai as genai
from pathlib import Path
from PIL import Image
from chemistry_tools import build_solver_hints
from chemistry_kb import retrieve_snippets

# Office file support
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("Warning: python-docx not available. Word document processing disabled.")

try:
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    print("Warning: openpyxl not available. Excel file processing disabled.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not available. PowerPoint file processing disabled.")

# PDF support
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("Warning: PyPDF2 not available. PDF processing disabled.")

# Import config for API key and settings
try:
    from config import GEMINI_API_KEY, MODEL_NAME
except ImportError:
    GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
    MODEL_NAME = "gemini-3.1-flash-lite-preview"

# Try to import vision tool for pix2text OCR
try:
    import sys
    from pathlib import Path as PathlibPath
    backend_path = str(PathlibPath(__file__).parent.parent / "backend")
    if backend_path not in sys.path:
        sys.path.insert(0, backend_path)
    from tools.vision_tool import VisionTool
    VISION_TOOL_AVAILABLE = True
except ImportError:
    VISION_TOOL_AVAILABLE = False
    print("Warning: Vision tool (pix2text) not available. Image analysis may be limited.")


class ChemistryAIBrain:
    """Backend brain for AI Chemistry Chatbot using Google Gemini API - Accurate Thinking Mode."""
    
    def __init__(self, api_key: Optional[str] = None):
        """
        Initialize the Chemistry AI Brain.
        
        Args:
            api_key: Google Gemini API key. If None, will try to load from GEMINI_API_KEY env variable.
        """
        self.api_key = api_key or GEMINI_API_KEY
        
        if not self.api_key:
            raise ValueError(
                "GEMINI_API_KEY not provided. Please set it as an environment variable "
                "or pass it to the constructor. "
                "Get your key at: https://aistudio.google.com/app/apikey"
            )
        
        genai.configure(api_key=self.api_key)
        
        self.system_prompt = """You are an intelligent Chemistry AI Router and Orchestrator for a high school chemistry tutoring system. Your primary role is to analyze student questions and intelligently delegate to specialized chemistry tools, then synthesize their outputs into coherent educational responses.

## Available Tools
- **pubchem_tool**: Fetches chemical properties, molar mass, structural information
- **wolfram_tool**: Balances equations, performs calculations, solves math problems
- **vision_tool**: Extracts text/equations from images using OCR (pix2text/Gemini Vision)
- **media_tool**: Generates educational images and audio content

## Your Workflow
1. **Analyze** the student's question to identify required tools
2. **Select** appropriate tools based on question type
3. **Orchestrate** tool calls in optimal sequence
4. **Synthesize** tool outputs into educational response
5. **Explain** results at high school level

## Tool Selection Logic
- Chemical properties/structures → pubchem_tool
- Equation balancing/calculations → wolfram_tool  
- Image analysis/diagrams → vision_tool
- Visual explanations needed → media_tool
- Complex problems → Multiple tools in sequence

## Response Format
Always structure responses with:
## Question Understanding
## Tools Used
## Step-by-Step Solution
## Key Concepts
## Practice Problem
## Summary

## Educational Approach (High School Focus)
- **Language**: Use Vietnamese when interacting with Vietnamese-speaking students
- **Vocabulary**: Maintain high school level - avoid quantum mechanics, advanced thermodynamics
- **Real-world connections**: Link concepts to everyday life (cooking, cleaning, environment)
- **Scaffolding**: Build from simple to complex, use analogies and examples
- **Confidence levels**: Indicate certainty for complex problems (High/Medium/Low)
- **Follow-up activities**: Suggest experiments, practice problems, or further reading

## High School Teaching Strategies
- **Conceptual first**: Explain "why" before "how"
- **Visual learning**: Use step-by-step breakdowns with clear formatting
- **Common misconceptions**: Address typical student errors explicitly
- **Study tips**: Include memory aids and learning strategies
- **Assessment prep**: Connect to exam formats and question types

## Response Structure for High School
1. **Hook**: Relatable example or question to engage interest
2. **Core concept**: Clear definition with simple explanation  
3. **Step-by-step**: Worked examples with reasoning
4. **Common mistakes**: What to avoid and why
5. **Practice**: Similar problem for student to try
6. **Connection**: Real-world application or career link

## Error Handling
If tools fail, provide conceptual explanations and suggest alternative approaches. Always prioritize educational value over perfect technical accuracy. Never leave a student without a learning path."""

        self.model = genai.GenerativeModel(
            MODEL_NAME,
            system_instruction=self.system_prompt
        )
        
        # Vision tool initialized lazily to avoid startup delays
        self._vision_tool = None
        self._vision_tool_initialized = False

    def _get_vision_tool(self):
        """Get vision tool, initializing lazily only when needed."""
        if not self._vision_tool_initialized:
            self._vision_tool_initialized = True
            if VISION_TOOL_AVAILABLE:
                try:
                    self._vision_tool = VisionTool(gemini_api_key=self.api_key)
                except Exception as e:
                    print(f"Warning: Could not initialize vision tool: {e}")
                    self._vision_tool = None
        return self._vision_tool
    
    @property
    def vision_tool(self):
        """Lazy property for vision tool."""
        return self._get_vision_tool()

    def _build_context_instruction(
        self,
        user_message: str,
        chemistry_context: Optional[Dict[str, str]] = None
    ) -> str:
        """Build tool orchestration context for each question."""
        context = chemistry_context or {}
        level = context.get("level", "highschool")
        style = context.get("style", "balanced")
        hints = build_solver_hints(user_message)
        snippets = retrieve_snippets(user_message, top_k=2)

        # Tool-specific instructions based on question analysis
        tool_instructions = self._analyze_and_select_tools(user_message)
        
        # Educational level mapping for high school focus
        level_map = {
            "highschool": "Use high-school level vocabulary. Avoid quantum mechanics. Focus on conceptual understanding.",
            "undergrad": "Use undergraduate-level depth with mechanisms and thermodynamics.",
            "mixed": "Provide layered explanation from intuitive to technical."
        }
        
        # Response style for educational effectiveness
        style_map = {
            "concise": "Direct answers with key equations and units.",
            "balanced": "Mix of explanation and step-by-step solutions.",
            "detailed": "Full derivations with concept connections."
        }

        snippet_lines = []
        for s in snippets:
            snippet_lines.append(f"- [{s['id']}] {s['content']} (source: {s['source']})")

        parts = [
            "## TOOL ORCHESTRATION INSTRUCTIONS",
            tool_instructions,
            "",
            "## EDUCATIONAL CONTEXT",
            f"Level: {level}. {level_map.get(level, level_map['highschool'])}",
            f"Style: {style}. {style_map.get(style, style_map['balanced'])}",
            "",
            "## CHEMISTRY HINTS",
            *[f"- {h}" for h in hints],
        ]
        
        if snippet_lines:
            parts.extend(["", "## KNOWLEDGE REFERENCES", *snippet_lines])
            
        return "\n".join(parts)

    def _analyze_and_select_tools(self, user_message: str) -> str:
        """Analyze question and determine which tools to use."""
        message_lower = user_message.lower()
        
        # Tool selection logic
        tools_needed = []
        
        # PubChem tool triggers
        if any(keyword in message_lower for keyword in ['molar mass', 'properties', 'structure', 'formula', 'element', 'compound']):
            tools_needed.append("pubchem_tool: Get chemical properties and structural data")
            
        # Wolfram tool triggers  
        if any(keyword in message_lower for keyword in ['balance', 'calculate', 'equation', 'math', 'solve', 'stoichiometry']):
            tools_needed.append("wolfram_tool: Balance equations and perform calculations")
            
        # Vision tool triggers
        if any(keyword in message_lower for keyword in ['image', 'diagram', 'picture', 'photo', 'ocr', 'read']):
            tools_needed.append("vision_tool: Extract text/equations from images")
            
        # Media tool triggers
        if any(keyword in message_lower for keyword in ['visualize', 'show', 'image', 'diagram', 'audio']):
            tools_needed.append("media_tool: Generate educational visuals/audio")
        
        if not tools_needed:
            tools_needed.append("Direct explanation: Use general chemistry knowledge")
            
        return "Tools to use:\n" + "\n".join(f"- {tool}" for tool in tools_needed)

    @staticmethod
    def _format_error(prefix: str, error: Exception) -> str:
        """Return a user-friendly API error message."""
        message = str(error)
        lowered = message.lower()
        if "401" in message or "unauthorized" in lowered:
            return (
                "Error: Unauthorized API request (401). "
                "Please check that `GEMINI_API_KEY` is set correctly and still valid."
            )
        return f"{prefix}: {message}"

    def _handle_tool_failure(self, tool_name: str, error: Exception, fallback_context: str) -> str:
        """Handle tool failures with educational fallbacks."""
        error_message = str(error).lower()
        
        # Provide educational fallback based on tool type
        if "pubchem" in tool_name.lower():
            return (
                f"⚠️ **PubChem Tool Unavailable**\n\n"
                f"Unable to fetch chemical data right now. "
                f"Here's what I can tell you from general chemistry knowledge:\n\n"
                f"{fallback_context}\n\n"
                f"💡 **Tip**: You can look up chemical properties using reliable sources like "
                f"PubChem directly or your chemistry textbook."
            )
        
        elif "wolfram" in tool_name.lower():
            return (
                f"⚠️ **Calculation Tool Unavailable**\n\n"
                f"Unable to perform calculations right now. "
                f"Let me guide you through the manual approach:\n\n"
                f"{fallback_context}\n\n"
                f"📝 **Manual Steps**: Show your work step-by-step using the formulas we've learned."
            )
        
        elif "vision" in tool_name.lower():
            return (
                f"⚠️ **Image Analysis Unavailable**\n\n"
                f"Cannot process images at the moment. "
                f"Please describe what you see in the image, and I'll help you analyze it.\n\n"
                f"📷 **Alternative**: You can type out equations or describe diagrams verbally."
            )
        
        elif "media" in tool_name.lower():
            return (
                f"⚠️ **Media Generation Unavailable**\n\n"
                f"Cannot generate visuals right now. "
                f"Let me explain the concept using text:\n\n"
                f"{fallback_context}\n\n"
                f"🎨 **Visualization Tip**: Try drawing the concept yourself based on the description."
            )
        
        else:
            return (
                f"⚠️ **Tool Unavailable**\n\n"
                f"Experiencing technical difficulties. "
                f"Here's a conceptual explanation:\n\n"
                f"{fallback_context}\n\n"
                f"🔄 **Try Again**: The issue might be temporary. Please retry in a moment."
            )

    def chat(
        self,
        user_message: str,
        chat_history: Optional[list] = None,
        chemistry_context: Optional[Dict[str, str]] = None,
        stream: bool = False
    ) -> Union[str, Iterator]:
        """
        Send a text message and get AI response (optionally streaming).
        
        Args:
            user_message: The user's chemistry question or message
            chat_history: Previous messages for context
            chemistry_context: Context for the model
            stream: Whether to stream the response
        
        Returns:
            AI's response text or a generator of response chunks
        """
        try:
            # Start a new chat session and build history
            chat_session = self.model.start_chat(history=[])
            
            if chat_history:
                history_for_api = []
                for msg in chat_history:
                    role = msg.get("role", "user")
                    if role == "assistant":
                        role = "model"
                    elif role not in {"user", "model"}:
                        role = "user"
                    history_for_api.append({"role": role, "parts": [msg.get("content", "")]})
                chat_session.history = history_for_api

            context_instruction = self._build_context_instruction(user_message, chemistry_context)
            enhanced_message = f"{context_instruction}\n\nUser question:\n{user_message}"

            # Get response from Gemini
            if stream:
                def response_generator():
                    response = chat_session.send_message(enhanced_message, stream=True)
                    for chunk in response:
                        if chunk.text:
                            yield chunk.text
                return response_generator()
            else:
                response = chat_session.send_message(enhanced_message)
                return response.text
        
        except Exception as e:
            return self._format_error("Error processing your question", e)

    def chat_with_image(
        self, 
        user_message: str, 
        image_path: str, 
        chat_history: Optional[list] = None,
        chemistry_context: Optional[Dict[str, str]] = None,
        stream: bool = False
    ) -> Union[str, Iterator]:
        """
        Send a message with an image and get AI response (optionally streaming).
        """
        try:
            if not Path(image_path).exists():
                return f"Error: Image file not found at {image_path}"
            
            # Extract text from image using pix2text (SimplePix2Text)
            extracted_content = ""
            vision_tool = self.vision_tool
            if vision_tool:
                try:
                    ocr_result = vision_tool.extract_from_image(image_path)
                    if ocr_result and ocr_result.success:
                        extracted_content = f"""📋 **Text/Equations Extracted from Image:**
- Extracted Text: {ocr_result.extracted_text}
- Equations Found: {', '.join(ocr_result.equations) if ocr_result.equations else 'None'}
- Chemical Formulas: {', '.join(ocr_result.chemical_formulas) if ocr_result.chemical_formulas else 'None'}
- Confidence: {ocr_result.confidence:.1%}

"""
                except Exception as e:
                    print(f"⚠️  pix2text extraction error: {str(e)}")
            
            # Load the image for Gemini analysis
            img = Image.open(image_path)
            
            # Build context instruction with extraction info
            context_instruction = self._build_context_instruction(
                f"{extracted_content}\n{user_message}" if extracted_content else user_message, 
                chemistry_context
            )
            
            # Send to Gemini
            prompt = [
                context_instruction, 
                f"{extracted_content}\n{user_message}" if extracted_content else user_message, 
                img
            ]
            
            if stream:
                def response_generator():
                    response = self.model.generate_content(prompt, stream=True)
                    for chunk in response:
                        if chunk.text:
                            yield chunk.text
                return response_generator()
            else:
                response = self.model.generate_content(prompt)
                return response.text
        
        except Exception as e:
            return self._format_error("Error processing image", e)

    def analyze_chemistry_concept(self, concept: str) -> str:
        """
        Get a detailed explanation of a chemistry concept.
        
        Args:
            concept: The chemistry concept to explain
        
        Returns:
            Detailed explanation
        """
        prompt = f"""Please provide a comprehensive explanation of the chemistry concept: "{concept}"

Include:
1. Definition
2. Key characteristics
3. Real-world applications
4. Examples
5. Related concepts
6. Common misconceptions (if any)"""
        
        return self.chat(prompt)

    def solve_problem(self, problem: str) -> str:
        """
        Solve a chemistry problem step-by-step.
        
        Args:
            problem: The chemistry problem to solve
        
        Returns:
            Step-by-step solution
        """
        prompt = f"""Please solve this chemistry problem step-by-step:

{problem}

Format your answer with:
1. Given information
2. What we need to find
3. Relevant formulas/equations
4. Step-by-step calculation
5. Final answer with units
6. Brief explanation of the result"""
        
        return self.chat(prompt)

    def get_molecule_info(self, molecule_name: str) -> str:
        """
        Get information about a molecule.
        
        Args:
            molecule_name: Name of the molecule
        
        Returns:
            Detailed molecule information
        """
        prompt = f"""Provide detailed information about {molecule_name}:

1. Chemical formula
2. Molecular structure (description)
3. Molar mass
4. Uses and applications
5. Safety information
6. Interesting facts"""
        
        return self.chat(prompt)

    def chat_with_file(
        self,
        user_message: str,
        file_path: str,
        chat_history: Optional[list] = None,
        chemistry_context: Optional[Dict[str, str]] = None,
        stream: bool = False
    ) -> Union[str, Iterator]:
        """
        Send a message with a file and get AI response (optionally streaming).
        Uses a combination of native File API and manual extraction for best results.
        """
        try:
            if not Path(file_path).exists():
                return f"Error: File not found at {file_path}"
            
            file_ext = Path(file_path).suffix.lower()
            
            # Handle images
            if file_ext in [".jpg", ".jpeg", ".png", ".gif", ".webp"]:
                return self.chat_with_image(user_message, file_path, chat_history, chemistry_context, stream=stream)
            
            # Handle text files
            elif file_ext in [".txt"]:
                return self._handle_text_file(user_message, file_path, chat_history, chemistry_context, stream=stream)
            
            # Handle PDF (Gemini native support is strong for PDF)
            elif file_ext == ".pdf":
                return self._handle_genai_file(user_message, file_path, chat_history, chemistry_context, stream=stream)
            
            # Handle Word, Excel, PPT (Manual extraction to avoid Unsupported MIME Type error)
            elif file_ext in [".doc", ".docx"]:
                return self._handle_docx_file(user_message, file_path, chat_history, chemistry_context, stream=stream)
            elif file_ext in [".xls", ".xlsx"]:
                return self._handle_excel_file(user_message, file_path, chat_history, chemistry_context, stream=stream)
            elif file_ext in [".ppt", ".pptx"]:
                return self._handle_pptx_file(user_message, file_path, chat_history, chemistry_context, stream=stream)
            
            else:
                return f"Error: Unsupported file type {file_ext}. Supported: jpg, png, gif, webp, txt, pdf, doc, docx, xls, xlsx, ppt, pptx"
        
        except Exception as e:
            return self._format_error("Error processing file", e)

    def _handle_genai_file(
        self,
        user_message: str,
        file_path: str,
        chat_history: Optional[list] = None,
        chemistry_context: Optional[Dict[str, str]] = None,
        stream: bool = False
    ) -> Union[str, Iterator]:
        """Handle document analysis using Gemini's native File API (best for PDFs)."""
        try:
            # Upload the file to Gemini
            genai_file = genai.upload_file(path=file_path, display_name=Path(file_path).name)
            
            context_instruction = self._build_context_instruction(user_message, chemistry_context)
            prompt = [
                context_instruction,
                f"The user has uploaded the file '{genai_file.display_name}'. Please analyze it and answer the following question.",
                genai_file,
                user_message
            ]

            if stream:
                def response_generator():
                    response = self.model.generate_content(prompt, stream=True)
                    for chunk in response:
                        if chunk.text:
                            yield chunk.text
                return response_generator()
            else:
                response = self.model.generate_content(prompt)
                return response.text
        except Exception as e:
            return self._format_error("Error processing file with Gemini API", e)

    def _handle_docx_file(
        self,
        user_message: str,
        file_path: str,
        chat_history: Optional[list] = None,
        chemistry_context: Optional[Dict[str, str]] = None,
        stream: bool = False
    ) -> Union[str, Iterator]:
        """Exhaustive Word document analysis to ensure no content (like Question 1) is missed."""
        try:
            if not DOCX_AVAILABLE:
                return "Error: python-docx is not installed."
            
            from docx import Document
            doc = Document(file_path)
            content_parts = []
            
            # 1. Check all document elements in order
            for element in doc.element.body:
                # Handle paragraphs
                if element.tag.endswith('p'):
                    from docx.text.paragraph import Paragraph
                    para = Paragraph(element, doc)
                    text = para.text.strip()
                    if text:
                        # Check style for header formatting
                        style = para.style.name.lower()
                        if "heading" in style:
                            content_parts.append(f"### {text}")
                        else:
                            content_parts.append(text)
                
                # Handle tables
                elif element.tag.endswith('tbl'):
                    from docx.table import Table
                    table = Table(element, doc)
                    content_parts.append("\n| " + " | ".join(["---"] * len(table.columns)) + " |")
                    for i, row in enumerate(table.rows):
                        row_text = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                        content_parts.append("| " + " | ".join(row_text) + " |")
                        if i == 0:
                            content_parts.append("| " + " | ".join(["---"] * len(row_text)) + " |")
            
            content = "\n\n".join(content_parts)
            if len(content) > 15000:
                content = content[:15000] + "\n\n[Content truncated...]"
                
            enhanced_message = f"Document '{Path(file_path).name}' for analysis.\n\nContent:\n{content}\n\n{user_message}"
            return self.chat(enhanced_message, chat_history, chemistry_context, stream=stream)
        except Exception as e:
            return self._format_error("Error processing Word document", e)

    def _handle_excel_file(
        self,
        user_message: str,
        file_path: str,
        chat_history: Optional[list] = None,
        chemistry_context: Optional[Dict[str, str]] = None,
        stream: bool = False
    ) -> Union[str, Iterator]:
        """Excel analysis with Markdown formatting."""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, data_only=True)
            content_parts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                content_parts.append(f"## Sheet: {sheet}")
                rows = list(ws.iter_rows(values_only=True))
                for i, row in enumerate(rows):
                    row_text = [str(cell).strip() if cell is not None else "" for cell in row]
                    if any(row_text):
                        content_parts.append("| " + " | ".join(row_text) + " |")
                        if i == 0:
                            content_parts.append("| " + " | ".join(["---"] * len(row_text)) + " |")
            content = "\n\n".join(content_parts)
            enhanced_message = f"Excel content:\n{content}\n\n{user_message}"
            return self.chat(enhanced_message, chat_history, chemistry_context, stream=stream)
        except Exception as e:
            return self._format_error("Error processing Excel", e)

    def _handle_pptx_file(
        self,
        user_message: str,
        file_path: str,
        chat_history: Optional[list] = None,
        chemistry_context: Optional[Dict[str, str]] = None,
        stream: bool = False
    ) -> Union[str, Iterator]:
        """PowerPoint analysis with slide context."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            content_parts = []
            for slide_idx, slide in enumerate(prs.slides, 1):
                content_parts.append(f"## Slide {slide_idx}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text.strip())
            content = "\n\n".join(content_parts)
            enhanced_message = f"PPT content:\n{content}\n\n{user_message}"
            return self.chat(enhanced_message, chat_history, chemistry_context, stream=stream)
        except Exception as e:
            return self._format_error("Error processing PowerPoint", e)



    def _handle_text_file(
        self,
        user_message: str,
        file_path: str,
        chat_history: Optional[list] = None,
        chemistry_context: Optional[Dict[str, str]] = None,
        stream: bool = False
    ) -> Union[str, Iterator]:
        """Handle text file analysis."""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
            
            # Limit content length
            if len(content) > 8000:
                content = content[:8000] + "\n\n[Content truncated for length...]"
            
            enhanced_message = f"I'm sharing a text file '{Path(file_path).name}' for analysis.\n\nFile content:\n---\n{content}\n---\n\n{user_message}"
            
            return self.chat(enhanced_message, chat_history, chemistry_context, stream=stream)
        
        except Exception as e:
            return self._format_error("Error reading text file", e)




# Initialize and cache the brain instance
_brain_instance = None


def get_brain(api_key: Optional[str] = None) -> ChemistryAIBrain:
    """
    Get or create the Chemistry AI Brain instance (singleton pattern).
    
    Args:
        api_key: Google Gemini API key (optional)
    
    Returns:
        ChemistryAIBrain instance
    """
    global _brain_instance
    if _brain_instance is None:
        _brain_instance = ChemistryAIBrain(api_key=api_key)
    return _brain_instance
